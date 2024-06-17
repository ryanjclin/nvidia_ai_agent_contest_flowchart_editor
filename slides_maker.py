
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from FlagEmbedding import BGEM3FlagModel
import numpy as np
import json
import os

def block_shape_decision(dict_shape_description, block_description):
    shape_names = [attr for attr in dir(MSO_SHAPE) if not callable(getattr(MSO_SHAPE, attr)) and not attr.startswith("__")]
    shape_dict = {name.lower(): getattr(MSO_SHAPE, name) for name in shape_names}

    model = BGEM3FlagModel('BAAI/bge-m3', use_fp16=True) 

    emb_block_description = model.encode(block_description, max_length=32)['dense_vecs']

    highest_score = 0
    best_shape = ""
    for key, value in dict_shape_description.items():
        api_description = f"{key}: {value}"
        emb_api_description = model.encode(api_description, max_length=32)['dense_vecs']

        score = emb_api_description @ emb_block_description.T

        if score > highest_score:
            best_shape = key
            highest_score = score

    
    # extract keyword
    best_shape = best_shape.lower()
    best_shape = best_shape.replace("msoshape","")

    lower_to_original = {k.lower(): k for k in shape_dict.keys() if isinstance(k, str)}
    lower_to_original.update({str(v).lower(): v for v in shape_dict['_valid_settings']})

    def search(query):
        query = query.lower().replace(' ', '')
        if query in lower_to_original:
            original_key = lower_to_original[query]
            if isinstance(original_key, str):
                return shape_dict[original_key]
            else:
                return shape_dict['_member_to_xml'][original_key]
        else:
            for k, v in shape_dict['_xml_to_member'].items():
                if k.lower().replace(' ', '') == query:
                    return v
        return None

    chosen_shape = search(best_shape)

    # print(f"block_description: {block_description}, chosen_shape: {chosen_shape} with score: {highest_score}")

    return chosen_shape

def creating_slides(json_data, resize_block = True):

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  
    slide = prs.slides.add_slide(slide_layout)

    ''' add blocks '''

    dict_block_info = {}

    ''' add block into slides'''
    for block in json_data['blocks']:
        if 'block' in block: block = block['block']

        # fetch data
        text_content = block['text content'][0]
        font_size = int(float(block['font size'][0]))
        font_color_r, font_color_b, font_color_g = int(block['text color'][0]), int(block['text color'][1]), int(block['text color'][2])
        block_color_r, block_color_b, block_color_g = int(block['block color'][0]), int(block['block color'][1]), int(block['block color'][2])
        center_x = int(float(block['Block center coordinates'][0]))
        center_y = int(float(block['Block center coordinates'][1]))
        block_width = int(float(block['block size'][0]))
        block_height = int(float(block['block size'][1]))
        block_description = block['block description'][0]

        # add ele into dict_block_info:
        dict_block_info[text_content] = [center_x, center_y, block_width, block_height]

        # coordinates_x, coordinates_y points to the top-left corner
        coordinates_x = int(center_x - (0.5 * block_width))
        coordinates_y = int(center_y - (0.5 * block_height))

        # determine block's shape
        json_file_path = 'clean_shape_description.json'
        with open(json_file_path, 'r') as file:
            dict_shape_description = json.load(file)

        dict_shape_description = json.load(open('clean_shape_description.json'))
        chosen_shape = block_shape_decision(dict_shape_description, block_description)

        coordinates_x_pt = Pt(coordinates_x)
        coordinates_y_pt = Pt(coordinates_y)
        block_width_pt = Pt(block_width) 
        block_height_pt = Pt(block_height)
        object_block = slide.shapes.add_shape(chosen_shape, coordinates_x_pt, coordinates_y_pt, block_width_pt, block_height_pt)

        # set content
        object_block.fill.solid()
        object_block.fill.fore_color.rgb = RGBColor(block_color_r, block_color_b, block_color_g)

        # font setting 
        text_frame = object_block.text_frame
        text_frame.text = text_content
        text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
        text_frame.paragraphs[0].font.size = Pt(font_size*0.85)



    ''' add arrows '''
    for arrow in json_data['arrows']:
        if 'arrow' in arrow: arrow = arrow['arrow']

        # fetch data
        starting_block_name = arrow['starting block'][0]
        ending_block_name = arrow['ending block'][0]
        starting_position_x_ai = int(float(arrow['start position'][0]))
        starting_position_y_ai = int(float(arrow['start position'][1]))
        ending_position_x_ai = int(float(arrow['end position'][0]))
        ending_position_y_ai = int(float(arrow['end position'][1]))
        arrow_color_r, arrow_color_b, arrow_color_g = int(arrow['color'][0]), int(arrow['color'][1]), int(arrow['color'][2])
        arrow_type = arrow['type'][0]
        arrow_boolean = arrow['arrow boolean'][0] if arrow['arrow boolean'][0] is not None else None

        # determine the trajectory of the arrow
        '''dict_block_info:
        key: text content
        value: center_x, center_y, block_width, block_height
        '''
        start_block_x = dict_block_info[starting_block_name][0]
        start_block_y = dict_block_info[starting_block_name][1]
        end_block_x = dict_block_info[ending_block_name][0]
        end_block_y = dict_block_info[ending_block_name][1]

        # horizon (y is the same)
        # start point x = the midpoint of right boundary of the start block
        # end point x = the midpoint of left boundary of the end block
        if end_block_y - 5 <= start_block_y and start_block_y <= end_block_y + 5:
            starting_position_x = int(start_block_x + (dict_block_info[starting_block_name][2] * 0.5)) 
            starting_position_y = start_block_y
            ending_position_x = int(end_block_x - (dict_block_info[ending_block_name][2] * 0.5)) 
            ending_position_y = end_block_y
        
        # vertical (x is the same)
        # start point y = the midpoint of buttom boundary of the start block
        # end point y = the midpoint of upper boundary of the end block
        elif end_block_x - 5 <= start_block_x and start_block_x <= end_block_x + 5:
            starting_position_x = start_block_x
            starting_position_y = int(start_block_y + (dict_block_info[starting_block_name][3] * 0.5)) 
            ending_position_x = end_block_x
            ending_position_y = int(end_block_y - (dict_block_info[ending_block_name][3] * 0.5)) 
            
        # determined by AI, not stable!!!
        else:   
            starting_position_x = starting_position_x_ai
            starting_position_y = starting_position_y_ai
            ending_position_x = ending_position_x_ai
            ending_position_y = ending_position_y_ai

        object_arrow = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(starting_position_x), Pt(starting_position_y), 
                                                Pt(ending_position_x), Pt(ending_position_y))
        object_arrow.line.color.rgb = RGBColor(arrow_color_r, arrow_color_b, arrow_color_g)
        object_arrow.line.width = Pt(2)
        # object_arrow.line.dash_style = MSO_LINE_DASH_STYLE.DASH 

        # add arrow head 
        connector_elm = object_arrow.element
        ln = connector_elm.find(".//a:ln", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
        if ln is not None:
            tail_end = parse_xml(f'<a:tailEnd {nsdecls("a")} type="triangle"/>')
            ln.append(tail_end)

        # add arrow boolean
        # if arrow_boolean is not None:
        if arrow_boolean != 'None':
            mid_x = Pt( (starting_position_x + ending_position_x) // 2 ) 
            mid_y = Pt( (starting_position_y + ending_position_y) // 2 )
            text_box = slide.shapes.add_textbox(mid_x + Pt(0.4), mid_y + Pt(0.4), Pt(0.2), Pt(0.2))
            text_box.text_frame.text = arrow_boolean
            text_box.line.color.rgb = RGBColor(0, 0, 0)


    file_path = './editable_flowchart_presentation.pptx'
    try:
        prs.save(file_path)
        print(f'Presentation saved successfully at {os.path.abspath(file_path)}')
    except Exception as e:
        print(f'Error saving presentation: {e}')

